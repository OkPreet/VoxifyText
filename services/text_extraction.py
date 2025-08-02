import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document

def extract_text_from_pdf(path):
    text = ""
    doc = fitz.open(path)
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_pptx(path):
    text = ""
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text(file_path, file_ext):
    if file_ext == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_ext == 'pptx':
        return extract_text_from_pptx(file_path)
    elif file_ext == 'docx':
        return extract_text_from_docx(file_path)
    elif file_ext == 'txt':
        return open(file_path, "r", encoding="utf-8").read()
    else:
        raise ValueError("Unsupported file type")
